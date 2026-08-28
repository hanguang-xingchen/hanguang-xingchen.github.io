#!/usr/bin/env python3
"""Extract text from .docx, .xlsx, .pptx files and output as plain text."""
import sys
import json
import os

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract table contents
    for i, table in enumerate(doc.tables):
        parts.append(f"\n--- Table {i+1} ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)

def extract_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True, read_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"\n=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=False):
            cells = []
            for cell in row:
                if cell.value is not None:
                    cells.append(f"{cell.coordinate}={cell.value}")
            if cells:
                parts.append("  " + ", ".join(cells))
    return "\n".join(parts)

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
            if shape.has_table:
                parts.append("  [Table]")
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append("  " + " | ".join(cells))
    return "\n".join(parts)

def main():
    if len(sys.argv) < 2:
        print("Usage: extract_office.py <file.docx|xlsx|pptx>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.isfile(path):
        print(f"Error: file not found: {path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            text = extract_docx(path)
        elif ext == ".xlsx":
            text = extract_xlsx(path)
        elif ext == ".pptx":
            text = extract_pptx(path)
        else:
            print(f"Error: unsupported format: {ext}", file=sys.stderr)
            sys.exit(1)
    except Exception as e:
        print(f"Error extracting {path}: {e}", file=sys.stderr)
        sys.exit(1)

    # Output as JSON so the model can parse structured content
    output = {
        "file": os.path.abspath(path),
        "format": ext[1:],  # docx, xlsx, pptx
        "text": text
    }
    print(json.dumps(output, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
